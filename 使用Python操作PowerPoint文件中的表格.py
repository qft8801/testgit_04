
#-*-coding:utf-8 -*-
import importlib
importlib.reload(sys)
import pptx
from pptx.util import Inches

#创建空白演示文档
pptFile = pptx.Presentation()
#插入一页幻灯片，插入表格
slide = pptFile.slides.add_slide(pptFile.slide_layouts[4])
tale = slide.shapes.add_table(rows =6,cols=4,left =Inches(1),top =Inches(2),width =Inches(8),height=Inches(4))
#遍历表格单元格，写入内容
for rowIndex,row in enumerate(table.table.rows):
	for colIndex, cell in enumerate(row.cells):
		if rowIndex==0:
			cell.text_frame.text = '列'+str(colIndex)
		else:
			cell.text_frame.text =str(rowIndex*colIndex)
		cell.margein_left =Inches(0.2)
pptFile.save('test3.pptx')

#打开已有演示文档，获取第一页幻灯片中的表格对象
pptFile =pptx.Presentation('test.pptx')
for shape in pptFile.slides[0].shapes:
	if	shape.shape_type==19:
		table = shape
		break

#遍历并输出单元格内容
for row in table.table.rows:
	for cell in row.cells:
		print(cell.text_frame.text,end='\t')
	print()
